import json
import re

from pptx import Presentation
from pptx.util import Inches, Pt


class SongHandler:
    def __init__(self, file_path, output_dir):
        self.file_path = file_path
        self.output_dir = output_dir
        self.presentation = Presentation(self.file_path)

    def parse_song(self):
        saved_song_dict = {}
        song_dict = {}
        song_list = []

        def is_empty_slide(slide):
            for shape in slide.shapes:
                if shape.text:
                    return False
            return True

        def add_song():
            nonlocal song_list, song_dict
            song_list.append(song_dict)
            saved_song_dict[title] = True
            song_dict = {}

        def parse_title(title):
            x = re.search(r"\(.*\)", title)
            if x != None:
                title = re.sub(r"\(.*\)", "", title)
                return f"{title.title().strip()} {x[0].upper()}"
            return title.title()

        def get_type_and_title(slide):
            title = slide.shapes[0].text.strip()
            title_sections = re.split(r"\s{2,}", title)
            if len(title_sections) == 1:
                return None, parse_title(title_sections[0])
            else:
                return title_sections[0].strip(), parse_title(title_sections[1])

        for slide_idx, slide in enumerate(self.presentation.slides):
            if len(slide.shapes) == 2 and not is_empty_slide(slide):
                type, title = get_type_and_title(slide)
                if title != song_dict.get("title"):
                    if song_dict.get("title") != None:
                        add_song()
                    if saved_song_dict.get(title) != None:
                        continue
                    song_dict["type"] = type
                    song_dict["title"] = title

                content = slide.shapes[1].text.strip().replace("\xa0", " ")
                if content != "":
                    if song_dict.get("content"):
                        song_dict["content"].append(content)
                    else:
                        song_dict["content"] = [content]
                if (
                    slide_idx == len(self.presentation.slides) - 1
                    and song_dict.get("title") != None
                ):
                    add_song()
            elif song_dict.get("title") != None:
                add_song()
        return song_list

    def convert_new(self):
        song_list = self.parse_song()
        self.write_json(song_list)

    def convert_to_16_9(self):
        left_margin = Inches(0.3)
        right_margin = Inches(0.3)

        for index, slide in enumerate(self.presentation.slides):
            all_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text
                    shape.left = left_margin
                    shape.width = (
                        self.presentation.slide_width - left_margin - right_margin
                    )

            if len(slide.shapes) == 2:
                body = slide.shapes[1]
                text_frame = body.text_frame

                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size in [Pt(54), Pt(48), Pt(44)]:
                            run.font.size = Pt(60)

            # empty slide if no text
            if all_text.strip() == "":
                self.__make_slide_empty(slide)

    def write_json(self, data):
        split = self.file_path.split("/")
        name = self.file_path.split("/")[len(split) - 1]
        # Serializing json
        json_object = json.dumps(data, indent=4, ensure_ascii=False)

        # Writing to sample.json
        with open(f"{self.output_dir}/song.json", "w", encoding="utf-8") as outfile:
            outfile.write(json_object)

    def save(self):
        split = self.file_path.split("/")
        name = self.file_path.split("/")[len(split) - 1]
        self.presentation.save(f"{self.output_dir}/{name}")

    def __make_slide_empty(self, slide):
        for shape in slide.shapes:
            slide.shapes._spTree.remove(shape._element)
