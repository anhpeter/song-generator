from pptx import Presentation


class SongEditor:

    def __init__(
        self,
        input_path: str,
    ):
        self.input_path = input_path
        self.template_presentation = Presentation(input_path)
        self.presentation = Presentation(input_path)

    def update_title_font_size(self, old_size, new_size):
        for slide_idx, slide in enumerate(self.presentation.slides):
            self.__update_title_font_size(slide, old_size, new_size)
            print(f"{round(((slide_idx+1)/len(self.presentation.slides)) * 100)}%")

    def __update_title_font_size(self, slide, old_size, new_size):
        for shape in slide.shapes:
            if shape.text:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size == old_size:
                            run.font.size = new_size

    def save(self, path):
        self.presentation.save(path)
