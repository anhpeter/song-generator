import copy

from pptx import Presentation


class SongGenerator:

    def __init__(
        self,
        input_template_file_path: str,
        input_song_list: list,
        content_max_length=28,
    ):
        self.content_max_length = content_max_length
        self.input_template_file_path = input_template_file_path
        self.input_song_list = input_song_list
        self.template_presentation = Presentation(input_template_file_path)
        self.presentation = Presentation(input_template_file_path)

    def duplicate_slide(self, pres, index):
        template = pres.slides[index]
        try:
            blank_slide_layout = pres.slide_layouts[index]
        except:
            blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)]

        copied_slide = pres.slides.add_slide(blank_slide_layout)

        for shp in template.shapes:
            el = shp.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, "p:extLs")

        # remove empty shape
        empty_shape_idx_list = []
        for idx, shape in enumerate(copied_slide.shapes):
            if shape.text.strip() == "":
                empty_shape_idx_list.append(idx)

        deleted_shape_count = 0
        for idx in empty_shape_idx_list:
            copied_slide.shapes.element.remove(
                copied_slide.shapes[idx - deleted_shape_count].element
            )
            deleted_shape_count += 1
        return copied_slide

    def delete_paragraph(self, paragraph):
        p = paragraph._p
        parent_element = p.getparent()
        parent_element.remove(p)

    def delete_run(self, run):
        r = run._r
        r.getparent().remove(r)

    def update_placeholder_content(self, slide, placeholder, content, title):
        try:
            for shape in slide.shapes:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text == f"<{placeholder}>":
                            if content == "" or content == None:
                                self.delete_run(run)
                            else:
                                run.text = content
                            return
        except Exception as e:
            print(f"Error update_placeholder_content: title:{title}; content:{content}")

    def new_slide(self):
        slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slide_layout)
        for shape in slide.shapes:
            slide.shapes._spTree.remove(shape._element)
        return slide

    def generate(self):
        slide_idx = 3
        for input_song in self.input_song_list:
            slide_idx += 1
            title_slide = self.duplicate_slide(self.presentation, 0)
            self.update_placeholder_content(
                title_slide, "type", input_song["type"], input_song["title"]
            )
            self.update_placeholder_content(
                title_slide, "title", input_song["title"], input_song["title"]
            )
            start_idx = 0
            for input_content in input_song["content"]:
                while True:
                    slide_content = input_content[start_idx:]
                    if len(slide_content.split(" ")) > 25:
                        end_idx = (
                            len(slide_content)
                            if len(slide_content) < self.content_max_length
                            else self.content_max_length
                        )
                        if self.content_max_length < len(slide_content):
                            last_c = slide_content[self.content_max_length]
                            if last_c != " ":
                                tmp: str = slide_content[:end_idx][::-1]
                                end_idx = self.content_max_length - tmp.index(" ")

                        slide_content = slide_content[:end_idx]
                    print(
                        f"slide:{slide_idx}:{len(slide_content.split(' '))}:{len(slide_content)}"
                    )
                    content_slide = self.duplicate_slide(self.presentation, 1)
                    self.update_placeholder_content(
                        content_slide,
                        "content",
                        slide_content.strip(),
                        input_song["title"],
                    )
                    start_idx = start_idx + len(slide_content)
                    slide_idx += 1
                    if start_idx >= len(input_content):
                        break
                start_idx = 0

    def save(self, path):
        self.presentation.save(path)
